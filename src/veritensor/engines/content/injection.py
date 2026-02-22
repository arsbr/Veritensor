# Copyright 2025 Veritensor Security Apache 2.0
# RAG Scanner: Detects Prompt Injections, PII, and Stealth Attacks (CSS/HTML hiding).

import logging
import re 
from typing import List, Generator, Set
from pathlib import Path
from veritensor.engines.static.rules import SignatureLoader
from veritensor.engines.content.pii import PIIScanner
from veritensor.core.text_utils import normalize_text, extract_base64_content
from veritensor.core.file_utils import validate_file_extension 
from veritensor.engines.content.injection import scan_text

logger = logging.getLogger(__name__)

# Supported text formats for RAG scanning
TEXT_EXTENSIONS = {
    # Documentation & Markup
    ".txt", ".md", ".markdown", ".rst", ".adoc", ".asciidoc", 
    ".tex", ".org", ".wiki", ".html", ".htm", ".css",
    
    # Data & Configs
    ".json", ".xml", ".yaml", ".yml", ".toml", 
    ".ini", ".cfg", ".conf", ".env", ".properties", ".editorconfig",
    
    # Source Code
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".hpp",
    ".rs", ".go", ".rb", ".php", ".pl", ".lua",
    ".sh", ".bash", ".zsh", ".ps1", ".bat", ".sql",
    
    # Infrastructure
    ".dockerfile", ".tf", ".tfvars", ".k8s", ".helm", ".tpl",
    ".gitignore", ".gitattributes",
    
    # Logs
    ".log", ".out", ".err"
}

DOC_EXTS = {".pdf", ".docx", ".pptx"}

# Import Optional Dependencies
try:
    import pypdf
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

CHUNK_SIZE = 1024 * 1024 # 1MB chunks
OVERLAP_SIZE = 4096      # 4KB overlap

# --- STEALTH ATTACK SIGNATURES (CSS/HTML Hiding) ---
STEALTH_PATTERNS = [
    r"font-size:\s*0px",
    r"font-size:\s*1px",
    r"color:\s*white",
    r"color:\s*#ffffff",
    r"color:\s*#fff",
    r"color:\s*transparent",
    r"display:\s*none",
    r"visibility:\s*hidden",
    r"opacity:\s*0",
    r"position:\s*absolute;\s*left:\s*-\d+px",
    r"z-index:\s*-\d+",
    r"<!--.*?ignore previous.*?-->", 
    r"<span[^>]*style=.*?>.*?</span>" 
]

# Copyright 2025 Veritensor Security Apache 2.0
# RAG Scanner: Detects Prompt Injections, PII, and Stealth Attacks (CSS/HTML hiding).

import logging
import re 
from typing import List, Generator, Set
from pathlib import Path
from veritensor.engines.static.rules import SignatureLoader
from veritensor.engines.content.pii import PIIScanner
# --- NEW IMPORT ---
from veritensor.core.text_utils import normalize_text, detect_stealth_text

logger = logging.getLogger(__name__)

# Supported text formats for RAG scanning
TEXT_EXTENSIONS = {
    # Documentation & Markup
    ".txt", ".md", ".markdown", ".rst", ".adoc", ".asciidoc", 
    ".tex", ".org", ".wiki", ".html", ".htm", ".css",
    
    # Data & Configs
    ".json", ".xml", ".yaml", ".yml", ".toml", 
    ".ini", ".cfg", ".conf", ".env", ".properties", ".editorconfig",
    
    # Source Code
    ".py", ".js", ".ts", ".java", ".c", ".cpp", ".h", ".hpp",
    ".rs", ".go", ".rb", ".php", ".pl", ".lua",
    ".sh", ".bash", ".zsh", ".ps1", ".bat", ".sql",
    
    # Infrastructure
    ".dockerfile", ".tf", ".tfvars", ".k8s", ".helm", ".tpl",
    ".gitignore", ".gitattributes",
    
    # Logs
    ".log", ".out", ".err"
}

DOC_EXTS = {".pdf", ".docx", ".pptx"}

# Import Optional Dependencies
try:
    import pypdf
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

CHUNK_SIZE = 1024 * 1024
OVERLAP_SIZE = 4096

STEALTH_PATTERNS = [
    r"font-size:\s*0px", r"font-size:\s*1px", r"color:\s*white", r"color:\s*#ffffff",
    r"color:\s*#fff", r"color:\s*transparent", r"display:\s*none", r"visibility:\s*hidden",
    r"opacity:\s*0", r"position:\s*absolute;\s*left:\s*-\d+px", r"z-index:\s*-\d+",
    r"<!--.*?ignore previous.*?-->", r"<span[^>]*style=.*?>.*?</span>" 
]

def scan_text(text: str, source_name: str = "memory") -> List[str]:
    """
    Core function to scan a raw string in memory.
    Used by File Scanners and In-Memory Integrations (Unstructured, ChromaDB).
    """
    threats = []
    if not text: return threats
    
    signatures = SignatureLoader.get_prompt_injections()
    
    # 1. Normalization
    clean_chunk = normalize_text(text)
    
    # 2. Base64 De-obfuscation
    decoded_parts = extract_base64_content(clean_chunk)
    for decoded in decoded_parts:
        for pattern in signatures:
            if is_match(decoded, [pattern]):
                 threats.append(f"HIGH: Obfuscated (Base64) Injection detected in {source_name}: '{pattern}'")

    # 3. Spaced/Obfuscated keywords
    compact_chunk = " ".join(clean_chunk.split())
    if len(compact_chunk) > 50 and (compact_chunk.count(" ") / len(compact_chunk)) > 0.3:
        collapsed_chunk = compact_chunk.replace(" ", "")
        CRITICAL_KEYWORDS = ["asananswer", "alwayswrite", "ignoreprevious", "systemoverride", "pwned", "jailbreak"]
        for kw in CRITICAL_KEYWORDS:
            if kw in collapsed_chunk.lower():
                threats.append(f"HIGH: Obfuscated/Spaced Injection detected in {source_name}: '{kw}'")
                return threats

    # 4. Standard Injection Logic
    for pattern in signatures:
        is_hit = False
        found_text = ""
        if pattern.startswith("regex:"):
            regex_str = pattern.replace("regex:", "", 1).strip()
            try:
                match = re.search(regex_str, compact_chunk, re.IGNORECASE)
                if match:
                    is_hit = True
                    found_text = match.group(0)[:100]
            except re.error:
                pass
        else:
            if pattern.lower() in compact_chunk.lower():
                is_hit = True
                found_text = pattern 
        
        if is_hit:
            threats.append(f"HIGH: Prompt Injection detected in {source_name}: Found '{found_text}'")
            return threats 

    # 5. PII Scan
    pii_threats = PIIScanner.scan(compact_chunk)
    if pii_threats:
        # Append source name to PII threats
        threats.extend([f"{t} in {source_name}" for t in pii_threats])
        return threats

    return threats

def scan_document(file_path: Path) -> List[str]:
    """Scans a file from disk."""
    ext = file_path.suffix.lower()
    threats = []
    
    try:
        # PHASE 0: Magic Numbers
        extension_threat = validate_file_extension(file_path)
        if extension_threat:
            return [extension_threat]

        # PHASE 1: Raw Content Scan (Stealth Detection)
        if ext in DOC_EXTS or ext in TEXT_EXTENSIONS:
            raw_threats = _scan_raw_binary(file_path)
            threats.extend(raw_threats)

        # PHASE 2: Extracted Text Scan
        text_generator = None
        if ext in TEXT_EXTENSIONS: text_generator = _read_text_sliding(file_path)
        elif ext == ".pdf" and PDF_AVAILABLE: text_generator = _yield_string_chunks(_read_pdf(file_path))
        elif ext == ".docx" and DOCX_AVAILABLE: text_generator = _yield_string_chunks(_read_docx(file_path))
        elif ext == ".pptx" and PPTX_AVAILABLE: text_generator = _yield_string_chunks(_extract_text_from_pptx(file_path))
        else: return threats 

        for chunk in text_generator:
            chunk_threats = scan_text(chunk, source_name=file_path.name)
            if chunk_threats:
                threats.extend(chunk_threats)
                return threats

    except Exception as e:
        logger.warning(f"Failed to scan document {file_path}: {e}")
        threats.append(f"WARNING: Doc Scan Error: {str(e)}")
        
    return threats

def _scan_raw_binary(path: Path) -> List[str]:
    threats = []
    try:
        with open(path, "rb") as f:
            buffer = ""
            while True:
                chunk_bytes = f.read(CHUNK_SIZE)
                if not chunk_bytes: break
                chunk_str = chunk_bytes.decode("latin-1")
                data = buffer + chunk_str
                for pattern in STEALTH_PATTERNS:
                    match = re.search(pattern, data, re.IGNORECASE)
                    if match:
                        found_text = match.group(0)
                        if len(found_text) > 50: found_text = found_text[:47] + "..."
                        threats.append(f"MEDIUM: Stealth/Hiding technique detected in {path.name} (Raw): '{found_text}'")
                        return threats 
                buffer = chunk_str[-OVERLAP_SIZE:]
    except Exception: pass 
    return threats

def _read_text_sliding(path: Path) -> Generator[str, None, None]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        buffer = ""
        while True:
            chunk = f.read(CHUNK_SIZE)
            if not chunk: break
            data = buffer + chunk
            yield data
            buffer = chunk[-OVERLAP_SIZE:]

def _yield_string_chunks(text: str) -> Generator[str, None, None]:
    if not text: return
    yield text

def _read_pdf(path: Path) -> str:
    text_content = []
    try:
        reader = pypdf.PdfReader(path)
        max_pages = min(len(reader.pages), 50) 
        for i in range(max_pages):
            page_text = reader.pages[i].extract_text()
            if page_text: text_content.append(page_text)
        return "\n".join(text_content)
    except Exception: return ""

def _read_docx(path: Path) -> str:
    text_content = []
    try:
        doc = docx.Document(path)
        max_paras = min(len(doc.paragraphs), 2000)
        for i in range(max_paras): text_content.append(doc.paragraphs[i].text)
        return "\n".join(text_content)
    except Exception: return ""

def _extract_text_from_pptx(path: Path) -> str:
    if not PPTX_AVAILABLE: return ""
    text_runs = []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): text_runs.append(shape.text)
    except Exception: pass
    return "\n".join(text_runs)



